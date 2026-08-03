import streamlit as st
import streamlit.components.v1 as components
import gspread
import base64
import json
import re
import time
import unicodedata
from datetime import datetime, date, timedelta
from google.oauth2 import service_account
from fpdf import FPDF

# Configuração da página - Centralizado e responsivo
st.set_page_config(
    page_title="Leitores Peregrinos",
    page_icon="https://i.ibb.co/j92LZnZJ/novo-logo-oficial.png",
    layout="centered",
    initial_sidebar_state="collapsed"
)

# Sobrescreve o ícone usado em "Adicionar à Tela de Início" (iOS e Android),
# que por padrão usa o manifesto/ícone genérico do Streamlit, não o nosso logo.
components.html("""
    <script>
        try {
            var doc = window.parent.document;
            var iconUrl = "https://i.ibb.co/j92LZnZJ/novo-logo-oficial.png";

            // iOS (Safari "Adicionar à Tela de Início")
            var appleLink = doc.querySelector('link[rel="apple-touch-icon"]');
            if (!appleLink) {
                appleLink = doc.createElement('link');
                appleLink.rel = 'apple-touch-icon';
                doc.head.appendChild(appleLink);
            }
            appleLink.href = iconUrl;

            // Android/Chrome (manifesto PWA)
            var manifestObj = {
                name: "Leitores Peregrinos",
                short_name: "Leitores Peregrinos",
                start_url: window.parent.location.href,
                display: "standalone",
                background_color: "#FEFAE0",
                theme_color: "#5C3A21",
                icons: [
                    { src: iconUrl, sizes: "192x192", type: "image/png" },
                    { src: iconUrl, sizes: "512x512", type: "image/png" }
                ]
            };
            var blob = new Blob([JSON.stringify(manifestObj)], {type: 'application/json'});
            var manifestUrl = URL.createObjectURL(blob);

            var manifestLink = doc.querySelector('link[rel="manifest"]');
            if (!manifestLink) {
                manifestLink = doc.createElement('link');
                manifestLink.rel = 'manifest';
                doc.head.appendChild(manifestLink);
            }
            manifestLink.href = manifestUrl;
        } catch (e) {
            console.log("[ICON-DEBUG] erro:", e);
        }
    </script>
""", height=0)

# --- GERENCIAMENTO DE ESTADO DA SESSÃO ---
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False
    st.session_state.user_name = ""
    st.session_state.user_profile = ""
    st.session_state.user_id = ""

if "pagina" not in st.session_state:
    st.session_state.pagina = "home"

if "ultima_pagina_rolada" not in st.session_state:
    st.session_state.ultima_pagina_rolada = None

# --- FUNÇÕES DE TRANSIÇÃO DE TELA (NATIVAS E RÁPIDAS) ---
def navegar_para(nome_pagina):
    st.session_state.pagina = nome_pagina

def efetuar_logout():
    st.session_state.logged_in = False
    st.session_state.user_name = ""
    st.session_state.user_profile = ""
    st.session_state.user_id = ""
    st.session_state.pagina = "home"

# --- FOLHA DE ESTILO GLOBAL (SEM PARAMETROS DE URL OU LINKS FALSOS) ---
st.markdown("""
    <style>
    /* Esconde a barra de ferramentas padrão do Streamlit (menu hambúrguer, rodapé, botão Deploy) */
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
    div[data-testid="stToolbar"] {visibility: hidden;}
    div[data-testid="stDecoration"] {visibility: hidden;}

    /* Forçar a largura ideal da página */
    .block-container {
        max-width: 700px !important;
        padding-top: 1.5rem !important;
        padding-bottom: 1.5rem !important;
    }
    
    /* Fundo Geral da Página (Marfim Quente) */
    .stApp {
        background-color: #FEFAE0 !important;
    }
    
    /* Ocultar elementos nativos do Streamlit */
    #MainMenu, footer, header {visibility: hidden !important;}
    
    /* Contêiner do Cabeçalho Superior Oficial */
    .cartao-superior-oficial {
        width: 100%;
        margin-bottom: 25px;
        display: flex;
        justify-content: center;
        align-items: center;
    }
    
    .imagem-layout-completo {
        width: 100% !important;
        height: auto !important;
        display: block;
        border-radius: 16px;
    }
    
    /* BARRA DE STATUS TERRACOTA UNIFICADA */
    .barra-status-alinhada {
        background-color: #EAB99F !important;
        border-radius: 20px !important;
        padding: 14px 20px !important;
        display: flex !important;
        justify-content: space-between !important;
        align-items: center !important;
        height: 56px !important;
        margin-bottom: 15px !important;
        box-shadow: 0 3px 6px rgba(0,0,0,0.05);
        width: 100%;
        box-sizing: border-box;
    }
    
    .texto-logado-interno {
        color: #3D2612 !important;
        font-weight: 700 !important;
        font-size: 16px !important;
        font-family: sans-serif;
    }

    /* POSICIONADOR DO BOTÃO SAIR NATIVO CHAVEADO (ao lado da barra, mesma altura) */
    .st-key-sair_wrapper {
        display: flex;
        align-items: center;
        justify-content: flex-end;
        height: 56px !important;
        margin-bottom: 15px !important;
    }

    .st-key-sair_wrapper button {
        background: #1C120C !important;
        color: #FFFFFF !important;
        border: 1px solid #3D2612 !important;
        border-radius: 25px !important;
        padding: 4px 24px !important;
        font-size: 15px !important;
        font-weight: bold !important;
        box-shadow: 0 4px 8px rgba(0,0,0,0.25) !important;
        height: 36px !important;
        width: auto !important;
    }
    
    .st-key-sair_wrapper button:hover {
        background: #2D1E15 !important;
    }
    
    /* ESTILIZAÇÃO DO PAINEL DO MENU (CINZA-GRAFITE) - AGORA EM COLUNA ÚNICA */
    .st-key-menu_grid {
        background-color: #4F5666 !important;
        border-radius: 16px !important;
        padding: 24px 16px !important;
        margin-bottom: 25px !important;
        box-shadow: inset 0 2px 6px rgba(0,0,0,0.2), 0 4px 10px rgba(0,0,0,0.1) !important;
    }
    
    /* RESOLUÇÃO VISUAL DOS BOTÕES INTERNOS (AZUL PETRÓLEO + DUPLO CONTORNO CHANFRADO BRONZE) */
    .stApp.stApp .st-key-menu_grid button,
    .stApp.stApp .st-key-menu_grid div.stLinkButton a {
        background: #0D1B2A !important;
        border: 3.5px solid #8C6D4F !important;
        outline: 1.5px solid #423224 !important;
        border-radius: 24px !important;
        padding: 12px 10px !important;
        font-size: 18px !important;
        font-weight: 700 !important;
        text-align: center !important;
        width: 100% !important;
        display: block !important;
        box-shadow: 0 4px 8px rgba(0,0,0,0.35), inset 0 1px 2px rgba(255,255,255,0.1) !important;
        font-family: sans-serif;
        text-decoration: none !important;
        box-sizing: border-box !important;
        height: auto !important;
        white-space: normal !important;
        line-height: 1.2 !important;
    }

    /* Cor do texto DENTRO do botão (o Streamlit envolve o rótulo do botão num stMarkdownContainer,
       então precisamos de mais especificidade do que a regra de texto de conteúdo abaixo) */
    .stApp.stApp .st-key-menu_grid button,
    .stApp.stApp .st-key-menu_grid button p,
    .stApp.stApp .st-key-menu_grid button div,
    .stApp.stApp .st-key-menu_grid button span,
    .stApp.stApp .st-key-menu_grid div.stLinkButton a,
    .stApp.stApp .st-key-menu_grid div.stLinkButton a p,
    .stApp.stApp .st-key-menu_grid div.stLinkButton a div,
    .stApp.stApp .st-key-menu_grid div.stLinkButton a span {
        color: #FFFFFF !important;
    }
    
    .stApp.stApp .st-key-menu_grid button:hover,
    .stApp.stApp .st-key-menu_grid div.stLinkButton a:hover {
        background: #15273C !important;
        border-color: #A38465 !important;
    }

    /* TÍTULO "Identificação do Leitor" (mesma paleta/estilo do cabeçalho oficial) */
    .titulo-identificacao {
        color: #5C3A21 !important;
        font-family: Georgia, 'Times New Roman', serif !important;
        font-weight: 700 !important;
        font-size: 34px !important;
        text-align: center;
        margin-bottom: 14px;
    }

    /* LABELS E TÍTULOS PADRÃO DO APP: sempre em tom escuro, mesmo se o celular estiver em modo escuro */
    .stApp h1, .stApp h2, .stApp h3, .stApp h4, .stApp h5, .stApp h6,
    .stApp label {
        color: #3D2612 !important;
    }

    /* ESPAÇAMENTO MAIS COMPACTO ENTRE OS EVENTOS (Escala Geral, Minha Escala, Aguardando Leitores) */
    .stApp hr {
        margin: 4px 0 !important;
    }
    div[data-testid="stElementContainer"] {
        margin-bottom: 0.2rem !important;
    }

    /* CAMPOS "Nome cadastrado" e "ID (Senha)" na tela de login: rótulo e texto digitado no mesmo tamanho do título */
    div[data-testid="stForm"] label {
        font-size: 22px !important;
        font-weight: 600 !important;
    }
    div[data-testid="stForm"] .stTextInput input {
        font-size: 22px !important;
        padding: 10px 12px !important;
    }

    /* TEXTO DE CONTEÚDO (descrições dos eventos em Escala Geral, Minha Escala, Aguardando Leitores etc.)
       Mira apenas o texto dentro de blocos de markdown/write - nunca o texto interno de botões. */
    .stApp div[data-testid="stMarkdownContainer"] p,
    .stApp div[data-testid="stMarkdownContainer"] li,
    .stApp div[data-testid="stMarkdownContainer"] {
        color: #3D2612 !important;
    }

    /* CAIXAS DE TEXTO (login e demais formulários): fundo branco, texto preto */
    .stTextInput input, .stTextArea textarea {
        background-color: #FFFFFF !important;
        color: #000000 !important;
        -webkit-text-fill-color: #000000 !important;
    }

    /* BOTÕES PADRÃO DO APP (Voltar ao Menu, Alterar, Salvar, Servir, Cancelar etc.) 
       Estes não têm key específico, então aplicamos um estilo legível (fundo escuro + texto dourado) */
    div[data-testid="stButton"] button {
        background-color: #0D1B2A !important;
        border: 1px solid #3D2612 !important;
        border-radius: 14px !important;
        font-weight: 600 !important;
    }
    .stApp.stApp div[data-testid="stButton"] button,
    .stApp.stApp div[data-testid="stButton"] button p,
    .stApp.stApp div[data-testid="stButton"] button div,
    .stApp.stApp div[data-testid="stButton"] button span {
        color: #FFFFFF !important;
    }
    div[data-testid="stButton"] button:hover {
        background-color: #15273C !important;
    }

    /* BOTÃO "Entrar" (login): no mesmo padrão visual dos botões do menu (pílula azul-petróleo + contorno bronze) */
    div[data-testid="stFormSubmitButton"] button {
        background-color: #0D1B2A !important;
        border: 3.5px solid #8C6D4F !important;
        outline: 1.5px solid #423224 !important;
        border-radius: 24px !important;
        padding: 12px 10px !important;
        font-size: 20px !important;
        font-weight: 700 !important;
        width: 100% !important;
        box-shadow: 0 4px 8px rgba(0,0,0,0.35), inset 0 1px 2px rgba(255,255,255,0.1) !important;
    }
    .stApp.stApp div[data-testid="stFormSubmitButton"] button,
    .stApp.stApp div[data-testid="stFormSubmitButton"] button p,
    .stApp.stApp div[data-testid="stFormSubmitButton"] button div,
    .stApp.stApp div[data-testid="stFormSubmitButton"] button span {
        color: #FFFFFF !important;
    }
    div[data-testid="stFormSubmitButton"] button:hover {
        background-color: #15273C !important;
        border-color: #A38465 !important;
    }

    /* O botão Sair e os botões do menu têm estilo próprio mais específico, então continuam intactos abaixo */
    .stApp.stApp .st-key-sair_wrapper button,
    .stApp.stApp .st-key-sair_wrapper button p,
    .stApp.stApp .st-key-sair_wrapper button div,
    .stApp.stApp .st-key-sair_wrapper button span {
        color: #FFFFFF !important;
    }
    </style>
""", unsafe_allow_html=True)

# --- CONEXÃO SEGURA COM O GOOGLE SHEETS ---
@st.cache_resource
def get_credentials():
    encoded_json = st.secrets["gcp_service_account"]["base64_json"]
    decoded_json = json.loads(base64.b64decode(encoded_json).decode('utf-8'))
    return service_account.Credentials.from_service_account_info(
        decoded_json,
        scopes=["https://www.googleapis.com/auth/spreadsheets"]
    )

def get_connection():
    creds = get_credentials()
    client = gspread.authorize(creds)
    return client.open_by_key("1RnwgFBWytspiM5eh5i0pgW2HXNRrwLXU4dYGXviDHlU")

@st.cache_data(ttl=60)
def carregar_dados_escala():
    try:
        sh = get_connection()
        ws_escala = sh.worksheet("Escala")
        return ws_escala.get_all_records()
    except Exception as e:
        st.error(f"Erro ao conectar com a base de dados: {e}")
        return []

@st.cache_data(ttl=60)
def obter_lista_leitores():
    try:
        sh = get_connection()
        ws_leitores = sh.worksheet("Nomes dos Leitores")
        dados = ws_leitores.get_all_values()
        leitores = []
        for row in dados[1:]:
            if len(row) > 0 and row[0].strip():
                leitores.append(row[0].strip())
        return sorted(leitores)
    except:
        return []

def normalizar_horario(txt):
    """Normaliza formatos de horário (10h, 10h00, 10:00, 10 horas...) para 'HH:MM'."""
    txt = str(txt).strip().lower().replace('horas', '').replace('h', ':').strip()
    txt = txt.rstrip(':').strip()
    if not txt:
        return ""
    partes = txt.split(':')
    try:
        h = int(partes[0].strip())
        m = int(partes[1].strip()) if len(partes) > 1 and partes[1].strip() != '' else 0
        return f"{h:02d}:{m:02d}"
    except (ValueError, IndexError):
        return txt

@st.cache_data(ttl=60)
def obter_horarios_padrao():
    """Lê a aba 'Horarios_Padrao' (TIPO_DIA | HORARIO) e agrupa os horários por tipo de dia."""
    try:
        sh = get_connection()
        ws_hp = sh.worksheet("Horarios_Padrao")
        dados = ws_hp.get_all_records()
        horarios_por_tipo = {}
        for r in dados:
            tipo = str(r.get('TIPO_DIA', '')).strip().upper()
            horario = normalizar_horario(r.get('HORARIO', ''))
            if tipo and horario:
                horarios_por_tipo.setdefault(tipo, []).append(horario)
        for tipo in horarios_por_tipo:
            horarios_por_tipo[tipo] = sorted(set(horarios_por_tipo[tipo]))
        return horarios_por_tipo
    except Exception:
        return {}

def horarios_disponiveis_para_data(data_selecionada, escala, horarios_padrao_data):
    """Lista fechada de horários possíveis para a data (sem digitação livre):
    1) Tenta pela tabela Horarios_Padrao (DOMINGO, SABADO, DIA_SEMANA, DIA4_SEMANA, DIA4_FDS);
    2) Se não achar, usa os horários já cadastrados na Escala para essa mesma data (ex: Solenidade)."""
    dia_semana = data_selecionada.weekday()  # Segunda=0 ... Domingo=6
    eh_fds = dia_semana in (5, 6)

    if data_selecionada.day == 4:
        tipo = "DIA4_FDS" if eh_fds else "DIA4_SEMANA"
    elif dia_semana == 6:
        tipo = "DOMINGO"
    elif dia_semana == 5:
        tipo = "SABADO"
    else:
        tipo = "DIA_SEMANA"

    lista = horarios_padrao_data.get(tipo, [])
    if lista:
        return lista

    horarios_reais = set()
    for r in escala:
        d = extrair_data_evento(str(r.get('DIA', '')))
        if d == data_selecionada:
            h = normalizar_horario(r.get('HORARIO', ''))
            if h:
                horarios_reais.add(h)
    return sorted(horarios_reais)

@st.cache_data(ttl=60)
def obter_roteiros():
    try:
        sh = get_connection()
        ws_roteiros = sh.worksheet("Roteiros")
        dados = ws_roteiros.get_all_records()
        roteiros = {}
        for r in dados:
            data_str = str(r.get('DIA', '')).strip()
            horario_str = normalizar_horario(r.get('HORARIO', ''))
            link = str(r.get('LINK', '')).strip()
            if data_str and link:
                roteiros[(data_str, horario_str)] = link
        return roteiros
    except Exception:
        return {}

def salvar_roteiro(sh, data_str, horario_str, link):
    ws_roteiros = sh.worksheet("Roteiros")
    dados = ws_roteiros.get_all_values()
    horario_norm = normalizar_horario(horario_str)
    for idx, row in enumerate(dados[1:], start=2):
        if len(row) > 1 and row[0].strip() == data_str and normalizar_horario(row[1]) == horario_norm:
            ws_roteiros.update_cell(idx, 3, link)
            return
    ws_roteiros.append_row([data_str, horario_str, link])

def excluir_roteiro(sh, data_str, horario_str):
    ws_roteiros = sh.worksheet("Roteiros")
    dados = ws_roteiros.get_all_values()
    horario_norm = normalizar_horario(horario_str)
    for idx, row in enumerate(dados[1:], start=2):
        if len(row) > 1 and row[0].strip() == data_str and normalizar_horario(row[1]) == horario_norm:
            ws_roteiros.delete_rows(idx)
            return True
    return False

MENSAGEM_PENALIDADE_PADRAO = "Atenção: você foi advertido(a) por uma falta. Na próxima falta sem aviso prévio, você ficará suspenso(a) por 30 dias."

@st.cache_data(ttl=60)
def obter_penalidades():
    try:
        sh = get_connection()
        ws_pen = sh.worksheet("Penalidades")
        dados = ws_pen.get_all_records()
        penalidades = {}
        for r in dados:
            leitor = str(r.get('LEITOR', '')).strip()
            mensagem = str(r.get('MENSAGEM', '')).strip()
            if leitor:
                penalidades[leitor.upper()] = mensagem or MENSAGEM_PENALIDADE_PADRAO
        return penalidades
    except Exception:
        return {}

def registrar_penalidade(sh, leitor):
    ws_pen = sh.worksheet("Penalidades")
    dados = ws_pen.get_all_values()
    for row in dados[1:]:
        if len(row) > 0 and row[0].strip().upper() == leitor.strip().upper():
            return  # já existe uma penalidade pendente para esse leitor
    ws_pen.append_row([leitor, MENSAGEM_PENALIDADE_PADRAO])

def consumir_penalidade(sh, leitor):
    """Remove a penalidade após ser exibida uma vez ao leitor; incrementa o Contador de Faltas
    e registra o timestamp do aviso na coluna Data-aviso, na aba 'Nomes dos Leitores'.
    Retorna a mensagem, se havia alguma."""
    ws_pen = sh.worksheet("Penalidades")
    dados = ws_pen.get_all_values()
    for idx, row in enumerate(dados[1:], start=2):
        if len(row) > 0 and row[0].strip().upper() == leitor.strip().upper():
            mensagem = row[1].strip() if len(row) > 1 and row[1].strip() else MENSAGEM_PENALIDADE_PADRAO
            ws_pen.delete_rows(idx)

            try:
                ws_leitores = sh.worksheet("Nomes dos Leitores")
                leitores_rows = ws_leitores.get_all_values()
                for r_idx, l_row in enumerate(leitores_rows[1:], start=2):
                    if len(l_row) > 0 and l_row[0].strip().upper() == leitor.strip().upper():
                        faltas_atual = l_row[2].strip() if len(l_row) > 2 and l_row[2].strip().isdigit() else "0"
                        novo_faltas = int(faltas_atual) + 1
                        ws_leitores.update_cell(r_idx, 3, str(novo_faltas))
                        timestamp_aviso = datetime.now().strftime("%d/%m/%Y %H:%M")
                        ws_leitores.update_cell(r_idx, 5, timestamp_aviso)
                        break
            except Exception:
                pass

            return mensagem
    return None

@st.cache_data(ttl=60)
def obter_suspensoes():
    try:
        sh = get_connection()
        ws_sus = sh.worksheet("Suspensoes")
        dados = ws_sus.get_all_records()
        suspensoes = {}
        for r in dados:
            leitor = str(r.get('LEITOR', '')).strip()
            data_inicio_str = str(r.get('DATA_INICIO', '')).strip()
            motivo = str(r.get('MOTIVO', '')).strip()
            if leitor and data_inicio_str:
                try:
                    data_inicio = datetime.strptime(data_inicio_str, "%d/%m/%Y").date()
                    data_fim = data_inicio + timedelta(days=30)
                    suspensoes[leitor.upper()] = {"data_fim": data_fim, "motivo": motivo}
                except ValueError:
                    pass
        return suspensoes
    except Exception:
        return {}

def registrar_suspensao(sh, leitor, motivo=""):
    ws_sus = sh.worksheet("Suspensoes")
    dados = ws_sus.get_all_values()
    hoje_str = date.today().strftime("%d/%m/%Y")
    for idx, row in enumerate(dados[1:], start=2):
        if len(row) > 0 and row[0].strip().upper() == leitor.strip().upper():
            ws_sus.update_cell(idx, 2, hoje_str)
            ws_sus.update_cell(idx, 3, motivo)
            return
    ws_sus.append_row([leitor, hoje_str, motivo])

def remover_suspensao(sh, leitor):
    ws_sus = sh.worksheet("Suspensoes")
    dados = ws_sus.get_all_values()
    for idx, row in enumerate(dados[1:], start=2):
        if len(row) > 0 and row[0].strip().upper() == leitor.strip().upper():
            ws_sus.delete_rows(idx)
            return True
    return False

def leitor_esta_suspenso(nome, suspensoes_data):
    """Retorna um dict {'data_fim': date, 'motivo': str} se o leitor estiver suspenso, ou None."""
    info = suspensoes_data.get(nome.strip().upper())
    if info and date.today() <= info["data_fim"]:
        return info
    return None

# --- FUNÇÕES DE VALIDAÇÃO E REGRAS DE NEGÓCIO ---
def usuario_ja_escalado_no_dia(escala, data_alvo, nome_usuario):
    for r in escala:
        if str(r.get('DIA', '')).strip() == str(data_alvo).strip():
            if (str(r.get('COMENTARISTA', '')).strip().upper() == nome_usuario.strip().upper() or 
                str(r.get('LEITURA1', '')).strip().upper() == nome_usuario.strip().upper() or 
                str(r.get('LEITURA2', '')).strip().upper() == nome_usuario.strip().upper()):
                return True
    return False

def contar_servicos_no_mes(escala, nome_usuario, data_referencia):
    """Conta quantas vezes nome_usuario já está escalado no MESMO MÊS/ANO de data_referencia."""
    count = 0
    for r in escala:
        data_evento = extrair_data_evento(str(r.get('DIA', '')))
        if data_evento is None or data_referencia is None:
            continue
        if data_evento.month != data_referencia.month or data_evento.year != data_referencia.year:
            continue
        if str(r.get('COMENTARISTA', '')).strip().upper() == nome_usuario.strip().upper():
            count += 1
        if str(r.get('LEITURA1', '')).strip().upper() == nome_usuario.strip().upper():
            count += 1
        if str(r.get('LEITURA2', '')).strip().upper() == nome_usuario.strip().upper():
            count += 1
    return count

def extrair_data_evento(dia_str):
    match = re.search(r'(\d{2}/\d{2}/\d{4})', str(dia_str))
    if match:
        try:
            return datetime.strptime(match.group(1), "%d/%m/%Y").date()
        except:
            pass
    return None

def processar_tentativa_cancelamento(sh, nome_usuario, dia_evento):
    data_evento = extrair_data_evento(dia_evento)
    if data_evento:
        hoje = date.today()
        vespera = data_evento - timedelta(days=1)
        if hoje >= vespera:
            st.error("Contate o coordenador para a sua substituição")
            try:
                ws_leitores = sh.worksheet("Nomes dos Leitores")
                leitores_rows = ws_leitores.get_all_values()
                for r_idx, l_row in enumerate(leitores_rows[1:], start=2):
                    if len(l_row) > 0 and l_row[0].strip().upper() == nome_usuario.strip().upper():
                        faltas_atual = l_row[2].strip() if len(l_row) > 2 and l_row[2].strip().isdigit() else "0"
                        novo_faltas = int(faltas_atual) + 1
                        ws_leitores.update_cell(r_idx, 3, str(novo_faltas))
                        break
            except Exception:
                pass
            return False
    return True

def eh_fim_de_semana(dia_str):
    d = dia_str.lower()
    return "sábado" in d or "sabado" in d or "domingo" in d

def eh_dia_4(dia_str):
    d = extrair_data_evento(dia_str)
    return d is not None and d.day == 4

def deve_exibir_comentarista(row):
    """Comentarista abre em fins de semana, solenidades, e todo dia 4 do mês (mesmo em dia de semana)."""
    dia = str(row.get('DIA', ''))
    solenidade = str(row.get('SOLENIDADE', 'NÃO')).strip().upper()
    return eh_fim_de_semana(dia) or solenidade == 'SIM' or eh_dia_4(dia)

def deve_exibir_leitura2(row):
    """2ª Leitura abre só em fins de semana e solenidades (dia 4 em dia de semana NÃO abre a 2ª Leitura)."""
    dia = str(row.get('DIA', ''))
    solenidade = str(row.get('SOLENIDADE', 'NÃO')).strip().upper()
    return eh_fim_de_semana(dia) or solenidade == 'SIM'

def evento_e_hoje_ou_futuro(row):
    data_evento = extrair_data_evento(str(row.get('DIA', '')))
    if data_evento is None:
        return True  # não foi possível identificar a data no texto; mantém o evento por segurança
    return data_evento >= date.today()

def data_valida_para_roteiro(data_roteiro, escala):
    if data_roteiro.weekday() in (5, 6):
        return True
    if data_roteiro.day == 4:
        return True
    for r in escala:
        dia_evento = extrair_data_evento(str(r.get('DIA', '')))
        solenidade = str(r.get('SOLENIDADE', 'NÃO')).strip().upper()
        if dia_evento == data_roteiro and solenidade == 'SIM':
            return True
    return False


# --- RENDERIZAÇÃO DO CABEÇALHO OFICIAL IMUTÁVEL ---
st.markdown("""
    <div class="cartao-superior-oficial">
        <img class="imagem-layout-completo" src="https://i.ibb.co/j92LZnZJ/novo-logo-oficial.png" />
    </div>
""", unsafe_allow_html=True)


# --- TELA DE LOGIN ---
if not st.session_state.logged_in:
    st.markdown('<div class="titulo-identificacao">Identificação do Leitor</div>', unsafe_allow_html=True)
    with st.form("login_form"):
        input_nome = st.text_input("Nome cadastrado:")
        input_senha = st.text_input("ID (Senha):", type="password")
        submit_login = st.form_submit_button("Entrar")
        
        if submit_login:
            try:
                sh = get_connection()
                ws_leitores = sh.worksheet("Nomes dos Leitores")
                leitores_data = ws_leitores.get_all_values()
                
                usuario_encontrado = False
                for idx, row in enumerate(leitores_data[1:], start=2):
                    if len(row) > 5 and row[0].strip().upper() == input_nome.strip().upper() and row[1].strip() == input_senha.strip():
                        st.session_state.logged_in = True
                        st.session_state.user_name = row[0].strip()
                        st.session_state.user_id = row[1].strip()
                        st.session_state.user_profile = row[5].strip()
                        usuario_encontrado = True
                        break
                
                if usuario_encontrado:
                    st.rerun()
                else:
                    st.error("Nome ou ID inválidos. Verifique seus dados.")
            except Exception as e:
                st.error(f"Erro ao conectar com a base de dados: {e}")
    st.stop()


# --- ÁREA LOGADA INTERNA ---
perfil_texto = "LEITOR"
if st.session_state.user_profile == "2":
    perfil_texto = "LEITOR & COMENTARISTA"
elif st.session_state.user_profile == "3":
    perfil_texto = "ADM"
elif st.session_state.user_profile == "5":
    perfil_texto = "ACOLHIDA"

# 1. Barra Terracota (identificação do Servo(a))
st.markdown(f"""
    <div class="barra-status-alinhada">
        <div class="texto-logado-interno">Servo(a): {st.session_state.user_name} ({perfil_texto})</div>
    </div>
""", unsafe_allow_html=True)

# 2. Instrução, logo abaixo da identificação do Servo(a), só na tela inicial
if st.session_state.pagina == "home":
    st.markdown('<div style="background-color: #A9ACB4; color: #10141A; border-radius: 8px; padding: 16px; text-align: center; font-size: 16px; font-weight: 600; margin-top: 10px; margin-bottom: 10px; font-family: sans-serif;">Selecione uma opção no menu abaixo para começar.</div>', unsafe_allow_html=True)

# 3. Menu de Botões Nativo, em coluna única (container com key="menu_grid" gera a classe .st-key-menu_grid usada no CSS acima)
# O botão "Sair" fica por último, dentro do próprio menu.
# Perfil "5" (Acolhida) só tem acesso a Coletar Intenções e Ver Intenções.
with st.container(key="menu_grid"):
    if st.session_state.user_profile == "5":
        st.link_button("Coletar Intenções", "https://docs.google.com/forms/d/e/1FAIpQLScgX8RkpDYhb-rMwb8_ZR6dJhp-tKUyowmRGrSK-tbsXveqCw/viewform?usp=sharing", use_container_width=True)
        st.button("Ver Intenções", key="menu_intencoes", on_click=navegar_para, args=("ver_intencoes",), use_container_width=True)
    else:
        st.button("Escala Geral", key="menu_geral", on_click=navegar_para, args=("escala_geral",), use_container_width=True)
        st.button("Minha Escala", key="menu_minha", on_click=navegar_para, args=("minha_escala",), use_container_width=True)
        st.button("Exibir Escala (PDF)", key="menu_pdf", on_click=navegar_para, args=("exibir_escala",), use_container_width=True)
        st.button("Aguardando Leitores", key="menu_vagas", on_click=navegar_para, args=("aguardando",), use_container_width=True)
        st.link_button("Coletar Intenções", "https://docs.google.com/forms/d/e/1FAIpQLScgX8RkpDYhb-rMwb8_ZR6dJhp-tKUyowmRGrSK-tbsXveqCw/viewform?usp=sharing", use_container_width=True)
        st.button("Ver Intenções", key="menu_intencoes", on_click=navegar_para, args=("ver_intencoes",), use_container_width=True)
        st.link_button("Liturgia Diária", "https://liturgia.cancaonova.com/pb/", use_container_width=True)
        if st.session_state.user_profile == "3":
            st.button("Roteiro", key="menu_roteiro", on_click=navegar_para, args=("cadastrar_roteiro",), use_container_width=True)
            st.button("Inserir Mensagem", key="menu_penalidade", on_click=navegar_para, args=("penalidade_leitor",), use_container_width=True)
            st.button("Suspensão de Leitor", key="menu_suspensao", on_click=navegar_para, args=("suspender_leitor",), use_container_width=True)
            st.button("Quem Não Serviu", key="menu_sem_escala", on_click=navegar_para, args=("quem_nao_serviu",), use_container_width=True)
    st.button("🚪 Sair", key="btn_logout_definitivo", on_click=efetuar_logout, use_container_width=True)

# Proteção extra: se por algum motivo o perfil 5 (Acolhida) estiver numa página que não é permitida
# para ele (ex: sessão antiga), volta para o menu principal.
if st.session_state.user_profile == "5" and st.session_state.pagina not in ("home", "ver_intencoes"):
    st.session_state.pagina = "home"


# Carregamento seguro dos dados da escala para os blocos abaixo
escala_data = carregar_dados_escala()
roteiros_data = obter_roteiros()
penalidades_data = obter_penalidades()
suspensoes_data = obter_suspensoes()
is_adm = (st.session_state.user_profile == "3")
lista_todos_leitores = obter_lista_leitores() if is_adm else []


# --- FUNÇÃO CENTRAL DE RENDERIZAÇÃO DOS EVENTOS ---
def renderizar_evento(idx, row, modo_aguardando=False):
    dia = str(row.get('DIA', ''))
    horario = str(row.get('HORARIO', ''))
    solenidade = str(row.get('SOLENIDADE', 'NÃO')).strip().upper()
    data_evento_atual = extrair_data_evento(dia)
    
    comentarista = str(row.get('COMENTARISTA', '')).strip()
    leitura1 = str(row.get('LEITURA1', '')).strip()
    leitura2 = str(row.get('LEITURA2', '')).strip()
    
    mostrar_comentarista = deve_exibir_comentarista(row)
    mostrar_l2 = deve_exibir_leitura2(row)
    
    tem_vago = (not leitura1) or (mostrar_comentarista and not comentarista) or (mostrar_l2 and not leitura2)
    if modo_aguardando and not tem_vago:
        return

    header_text = f"📅 **{dia}** — ⏰ **{horario}**"
    if solenidade == 'SIM':
        header_text += " ✨ *(Solenidade)*"
    if eh_dia_4(dia):
        header_text += " *(Missa da Saúde)*"
    st.markdown(header_text)

    if eh_dia_4(dia):
        st.markdown(
            '<img src="https://i.ibb.co/5XT3ZnxX/Whats-App-Image-2026-07-23-at-09-57-56.jpg" '
            'style="height:31px; width:auto; vertical-align:middle; margin:-8px 0 4px 0;">',
            unsafe_allow_html=True
        )

    data_evento = extrair_data_evento(dia)
    if data_evento:
        data_evento_fmt = data_evento.strftime("%d/%m/%Y")
        chave_roteiro = (data_evento_fmt, normalizar_horario(horario))
        if chave_roteiro in roteiros_data:
            st.markdown(f"📄 [Roteiro desta missa]({roteiros_data[chave_roteiro]})")
    
    usuario_atual = st.session_state.user_name

    if f"alterando_com_{idx}" not in st.session_state:
        st.session_state[f"alterando_com_{idx}"] = False
    if f"alterando_l1_{idx}" not in st.session_state:
        st.session_state[f"alterando_l1_{idx}"] = False
    if f"alterando_l2_{idx}" not in st.session_state:
        st.session_state[f"alterando_l2_{idx}"] = False

    # 1. COMENTARISTA
    if mostrar_comentarista:
        val_com = comentarista if comentarista else "Vago"
        c_col1, c_col2 = st.columns([3, 1.5])
        c_col1.write(f"**COMENTARISTA:** {val_com}")
        
        if is_adm:
            if not st.session_state[f"alterando_com_{idx}"]:
                if c_col2.button("Alterar", key=f"btn_alt_com_{idx}"):
                    st.session_state[f"alterando_com_{idx}"] = True
                    st.rerun()
            else:
                novo_nome_com = c_col2.selectbox("Novo Leitor:", ["(Vago)"] + lista_todos_leitores, key=f"sel_com_{idx}")
                if c_col2.button("Salvar", key=f"save_com_{idx}"):
                    val_salvar = "" if novo_nome_com == "(Vago)" else novo_nome_com
                    sh_conn = get_connection()
                    ws_live = sh_conn.worksheet("Escala")
                    ws_live.update_cell(idx + 2, 4, val_salvar)
                    carregar_dados_escala.clear()
                    st.session_state[f"alterando_com_{idx}"] = False
                    st.success("Alterado com sucesso!")
                    time.sleep(2.5)
                    st.rerun()
        else:
            if not comentarista:
                if c_col2.button("Servir", key=f"s_com_{idx}"):
                    info_suspensao = leitor_esta_suspenso(usuario_atual, suspensoes_data)
                    if info_suspensao:
                        msg_susp = f"Você está suspenso(a) até {info_suspensao['data_fim'].strftime('%d/%m/%Y')} e não pode se escalar."
                        if info_suspensao["motivo"]:
                            msg_susp += f" Motivo: {info_suspensao['motivo']}"
                        st.error(msg_susp)
                    elif st.session_state.user_profile == "1":
                        st.error("Você não possui o perfil “Comentarista”")
                    elif contar_servicos_no_mes(escala_data, usuario_atual, data_evento_atual) >= 3:
                        st.error("Você já serviu três vezes nesse mês")
                    elif usuario_ja_escalado_no_dia(escala_data, dia, usuario_atual):
                        st.error("Você já possui uma função agendada neste dia.")
                    else:
                        sh_conn = get_connection()
                        ws_live = sh_conn.worksheet("Escala")
                        ws_live.update_cell(idx + 2, 4, usuario_atual)
                        carregar_dados_escala.clear()
                        mensagem_penalidade = consumir_penalidade(sh_conn, usuario_atual)
                        st.success("Escalado como Comentarista!")
                        if mensagem_penalidade:
                            st.warning(mensagem_penalidade)
                        time.sleep(2.5)
                        st.rerun()
            else:
                if comentarista.upper() == usuario_atual.upper():
                    if c_col2.button("Cancelar", key=f"c_com_{idx}"):
                        sh_conn = get_connection()
                        if processar_tentativa_cancelamento(sh_conn, usuario_atual, dia):
                            ws_live = sh_conn.worksheet("Escala")
                            ws_live.update_cell(idx + 2, 4, "")
                            carregar_dados_escala.clear()
                            st.success("Cancelado com sucesso!")
                            time.sleep(2.5)
                            st.rerun()

    # 2. 1ª LEITURA
    val_l1 = leitura1 if leitura1 else "Vago"
    l1_col1, l1_col2 = st.columns([3, 1.5])
    l1_col1.write(f"**1ª LEITURA:** {val_l1}")

    if is_adm:
        if not st.session_state[f"alterando_l1_{idx}"]:
            if l1_col2.button("Alterar", key=f"btn_alt_l1_{idx}"):
                st.session_state[f"alterando_l1_{idx}"] = True
                st.rerun()
        else:
            novo_nome_l1 = l1_col2.selectbox("Novo Leitor:", ["(Vago)"] + lista_todos_leitores, key=f"sel_l1_{idx}")
            if l1_col2.button("Salvar", key=f"save_l1_{idx}"):
                val_salvar = "" if novo_nome_l1 == "(Vago)" else novo_nome_l1
                sh_conn = get_connection()
                ws_live = sh_conn.worksheet("Escala")
                ws_live.update_cell(idx + 2, 5, val_salvar)
                carregar_dados_escala.clear()
                st.session_state[f"alterando_l1_{idx}"] = False
                st.success("Alterado com sucesso!")
                time.sleep(2.5)
                st.rerun()
    else:
        if not leitura1:
            if l1_col2.button("Servir", key=f"s_l1_{idx}"):
                info_suspensao = leitor_esta_suspenso(usuario_atual, suspensoes_data)
                if info_suspensao:
                    msg_susp = f"Você está suspenso(a) até {info_suspensao['data_fim'].strftime('%d/%m/%Y')} e não pode se escalar."
                    if info_suspensao["motivo"]:
                        msg_susp += f" Motivo: {info_suspensao['motivo']}"
                    st.error(msg_susp)
                elif contar_servicos_no_mes(escala_data, usuario_atual, data_evento_atual) >= 3:
                    st.error("Você já serviu três vezes nesse mês")
                elif usuario_ja_escalado_no_dia(escala_data, dia, usuario_atual):
                    st.error("Você já possui uma função agendada neste dia.")
                else:
                    sh_conn = get_connection()
                    ws_live = sh_conn.worksheet("Escala")
                    ws_live.update_cell(idx + 2, 5, usuario_atual)
                    carregar_dados_escala.clear()
                    mensagem_penalidade = consumir_penalidade(sh_conn, usuario_atual)
                    st.success("Escalado na 1ª Leitura!")
                    if mensagem_penalidade:
                        st.warning(mensagem_penalidade)
                    time.sleep(2.5)
                    st.rerun()
        else:
            if leitura1.upper() == usuario_atual.upper():
                if l1_col2.button("Cancelar", key=f"c_l1_{idx}"):
                    sh_conn = get_connection()
                    if processar_tentativa_cancelamento(sh_conn, usuario_atual, dia):
                        ws_live = sh_conn.worksheet("Escala")
                        ws_live.update_cell(idx + 2, 5, "")
                        carregar_dados_escala.clear()
                        st.success("Cancelado com sucesso!")
                        time.sleep(2.5)
                        st.rerun()

    # 3. 2ª LEITURA
    if mostrar_l2:
        val_l2 = leitura2 if leitura2 else "Vago"
        l2_col1, l2_col2 = st.columns([3, 1.5])
        l2_col1.write(f"**2ª LEITURA:** {val_l2}")

        if is_adm:
            if not st.session_state[f"alterando_l2_{idx}"]:
                if l2_col2.button("Alterar", key=f"btn_alt_l2_{idx}"):
                    st.session_state[f"alterando_l2_{idx}"] = True
                    st.rerun()
            else:
                novo_nome_l2 = l2_col2.selectbox("Novo Leitor:", ["(Vago)"] + lista_todos_leitores, key=f"sel_l2_{idx}")
                if l2_col2.button("Salvar", key=f"save_l2_{idx}"):
                    val_salvar = "" if novo_nome_l2 == "(Vago)" else novo_nome_l2
                    sh_conn = get_connection()
                    ws_live = sh_conn.worksheet("Escala")
                    ws_live.update_cell(idx + 2, 6, val_salvar)
                    carregar_dados_escala.clear()
                    st.session_state[f"alterando_l2_{idx}"] = False
                    st.success("Alterado com sucesso!")
                    time.sleep(2.5)
                    st.rerun()
        else:
            if not leitura2:
                if l2_col2.button("Servir", key=f"s_l2_{idx}"):
                    info_suspensao = leitor_esta_suspenso(usuario_atual, suspensoes_data)
                    if info_suspensao:
                        msg_susp = f"Você está suspenso(a) até {info_suspensao['data_fim'].strftime('%d/%m/%Y')} e não pode se escalar."
                        if info_suspensao["motivo"]:
                            msg_susp += f" Motivo: {info_suspensao['motivo']}"
                        st.error(msg_susp)
                    elif contar_servicos_no_mes(escala_data, usuario_atual, data_evento_atual) >= 3:
                        st.error("Você já serviu três vezes nesse mês")
                    elif usuario_ja_escalado_no_dia(escala_data, dia, usuario_atual):
                        st.error("Você já possui uma função agendada neste dia.")
                    else:
                        sh_conn = get_connection()
                        ws_live = sh_conn.worksheet("Escala")
                        ws_live.update_cell(idx + 2, 6, usuario_atual)
                        carregar_dados_escala.clear()
                        mensagem_penalidade = consumir_penalidade(sh_conn, usuario_atual)
                        st.success("Escalado na 2ª Leitura!")
                        if mensagem_penalidade:
                            st.warning(mensagem_penalidade)
                        time.sleep(2.5)
                        st.rerun()
            else:
                if leitura2.upper() == usuario_atual.upper():
                    if l2_col2.button("Cancelar", key=f"c_l2_{idx}"):
                        sh_conn = get_connection()
                        if processar_tentativa_cancelamento(sh_conn, usuario_atual, dia):
                            ws_live = sh_conn.worksheet("Escala")
                            ws_live.update_cell(idx + 2, 6, "")
                            carregar_dados_escala.clear()
                            st.success("Cancelado com sucesso!")
                            time.sleep(2.5)
                            st.rerun()

    st.markdown("---")

# --- ROTEAMENTO E CONTEÚDOS ---
pagina_mudou = (st.session_state.pagina != st.session_state.ultima_pagina_rolada)

if st.session_state.pagina != "home":
    st.button("⬅️ Voltar ao Menu Principal", key="btn_voltar_menu", on_click=navegar_para, args=("home",))
    st.markdown('<div id="ancora-conteudo"></div>', unsafe_allow_html=True)

    # Rola automaticamente até o início do conteúdo desta tela - SÓ quando a página muda de verdade
    # (evita perder a posição de rolagem quando o ADM clica em Alterar/Salvar na mesma tela)
    if pagina_mudou:
        _nonce_scroll = f"{st.session_state.pagina}-{time.time()}"
        components.html(f"""
            <!-- nonce:{_nonce_scroll} -->
            <script>
                setTimeout(function() {{
                    try {{
                        var doc = window.parent.document;
                        var el = doc.getElementById("ancora-conteudo");
                        if (!el) {{ return; }}
                        var contentEl = el;
                        var scrollable = null;
                        while (contentEl && contentEl !== doc.body) {{
                            var style = doc.defaultView.getComputedStyle(contentEl);
                            if ((style.overflowY === "auto" || style.overflowY === "scroll") && contentEl.scrollHeight > contentEl.clientHeight) {{
                                scrollable = contentEl;
                                break;
                            }}
                            contentEl = contentEl.parentElement;
                        }}
                        if (scrollable) {{
                            var elRect = el.getBoundingClientRect();
                            var scrollableRect = scrollable.getBoundingClientRect();
                            var alvo = scrollable.scrollTop + (elRect.top - scrollableRect.top);
                            scrollable.scrollTop = alvo;
                        }} else {{
                            el.scrollIntoView({{behavior: "auto", block: "start"}});
                        }}
                    }} catch (e) {{
                        console.log("[SCROLL-DEBUG] ERRO:", e);
                    }}
                }}, 20);
            </script>
        """, height=0)
        st.session_state.ultima_pagina_rolada = st.session_state.pagina
else:
    # De volta ao menu principal: rola instantaneamente para o topo da página - só na primeira vez que chega no home
    if pagina_mudou:
        _nonce_topo = f"home-{time.time()}"
        components.html(f"""
            <!-- nonce:{_nonce_topo} -->
            <script>
                setTimeout(function() {{
                    try {{
                        var doc = window.parent.document;
                        var candidatos = doc.querySelectorAll('section, div, main');
                        candidatos.forEach(function(el) {{
                            var style = doc.defaultView.getComputedStyle(el);
                            if ((style.overflowY === "auto" || style.overflowY === "scroll") && el.scrollHeight > el.clientHeight) {{
                                el.scrollTop = 0;
                            }}
                        }});
                        doc.documentElement.scrollTop = 0;
                        doc.body.scrollTop = 0;
                        window.parent.scrollTo(0, 0);
                    }} catch (e) {{
                        console.log("[SCROLL-DEBUG] ERRO ao rolar para o topo:", e);
                    }}
                }}, 20);
            </script>
        """, height=0)
        st.session_state.ultima_pagina_rolada = "home"

if st.session_state.pagina == "home":
    pass

elif st.session_state.pagina == "cadastrar_roteiro":
    st.subheader("Cadastrar Roteiro")
    if st.session_state.user_profile != "3":
        st.error("Apenas o ADM pode acessar esta tela.")
    else:
        st.write("Selecione a data da missa (sábado, domingo, dia 4 do mês, ou dia marcado como Solenidade na Escala).")
        data_roteiro = st.date_input("Data da missa:", format="DD/MM/YYYY", key="data_roteiro_input")

        horarios_padrao_data = obter_horarios_padrao()
        opcoes_horario = horarios_disponiveis_para_data(data_roteiro, escala_data, horarios_padrao_data)

        if not opcoes_horario:
            st.warning("Não há horário de missa cadastrado para essa data (nem na tabela de horários padrão, nem na Escala Geral). Cadastre o horário na aba 'Horarios_Padrao' ou na 'Escala' antes de continuar.")
        else:
            horario_roteiro = st.selectbox("Horário da missa:", opcoes_horario, key="horario_roteiro_select")
            link_roteiro = st.text_input("Link do arquivo de roteiro:")

            if st.button("Salvar Roteiro"):
                if not data_valida_para_roteiro(data_roteiro, escala_data):
                    st.error("A data deve ser um sábado, domingo, dia 4 do mês, ou um dia marcado como Solenidade na Escala Geral.")
                elif not link_roteiro.strip():
                    st.error("Informe o link do arquivo de roteiro.")
                else:
                    sh_conn = get_connection()
                    data_str = data_roteiro.strftime("%d/%m/%Y")
                    salvar_roteiro(sh_conn, data_str, horario_roteiro, link_roteiro.strip())
                    obter_roteiros.clear()
                    st.success(f"Roteiro salvo para {data_str} às {horario_roteiro}!")
                    time.sleep(2.5)
                    st.rerun()

        st.markdown("---")
        st.write("**Roteiros já cadastrados:**")
        if not roteiros_data:
            st.info("Nenhum roteiro cadastrado ainda.")
        else:
            for (data_str, horario_str), link in sorted(roteiros_data.items()):
                col_link, col_excluir = st.columns([4, 1])
                col_link.markdown(f"- **{data_str} às {horario_str}**: [{link}]({link})")
                if col_excluir.button("🗑️ Excluir", key=f"excluir_roteiro_{data_str}_{horario_str}"):
                    sh_conn = get_connection()
                    if excluir_roteiro(sh_conn, data_str, horario_str):
                        obter_roteiros.clear()
                        st.success(f"Roteiro de {data_str} às {horario_str} removido!")
                        time.sleep(2.5)
                        st.rerun()
                    else:
                        st.error("Não foi possível encontrar esse roteiro para remover.")

elif st.session_state.pagina == "penalidade_leitor":
    st.subheader("Inserir Mensagem de Penalidade")
    if st.session_state.user_profile != "3":
        st.error("Apenas o ADM pode acessar esta tela.")
    else:
        st.write("Selecione o leitor que teve uma falta. Na próxima vez que ele(a) se escalar, receberá o aviso de que a próxima falta sem avisar resultará em suspensão de 30 dias.")

        if not lista_todos_leitores:
            st.info("Nenhum leitor cadastrado encontrado.")
        else:
            leitor_penalidade = st.selectbox("Selecione o Leitor:", lista_todos_leitores, key="sel_penalidade")

            if st.button("Registrar Penalidade"):
                sh_conn = get_connection()
                registrar_penalidade(sh_conn, leitor_penalidade)
                obter_penalidades.clear()
                st.success(f"Penalidade registrada para {leitor_penalidade}.")
                time.sleep(2.5)
                st.rerun()

            st.markdown("---")
            st.write("**Penalidades pendentes (ainda não exibidas ao leitor):**")
            if not penalidades_data:
                st.info("Nenhuma penalidade pendente no momento.")
            else:
                for leitor_nome in sorted(penalidades_data.keys()):
                    st.markdown(f"- **{leitor_nome}**")

elif st.session_state.pagina == "suspender_leitor":
    st.subheader("Suspensão de Leitor")
    if st.session_state.user_profile != "3":
        st.error("Apenas o ADM pode acessar esta tela.")
    else:
        st.write("Selecione o leitor a ser suspenso. A suspensão dura 30 dias corridos, contados a partir de hoje.")

        if not lista_todos_leitores:
            st.info("Nenhum leitor cadastrado encontrado.")
        else:
            leitor_suspender = st.selectbox("Selecione o Leitor:", lista_todos_leitores, key="sel_suspender")
            motivo_suspensao = st.text_area("Motivo da suspensão (será exibido ao leitor):", key="motivo_suspender", placeholder="Ex: Ausências recorrentes sem aviso prévio nos últimos meses.")

            if st.button("Suspender por 30 dias"):
                sh_conn = get_connection()
                registrar_suspensao(sh_conn, leitor_suspender, motivo_suspensao.strip())
                obter_suspensoes.clear()
                data_fim_prevista = date.today() + timedelta(days=30)
                st.success(f"{leitor_suspender} suspenso(a) até {data_fim_prevista.strftime('%d/%m/%Y')}.")
                time.sleep(2.5)
                st.rerun()

            st.markdown("---")
            st.write("**Leitores atualmente suspensos:**")
            suspensos_ativos = {nome: info for nome, info in suspensoes_data.items() if info["data_fim"] >= date.today()}
            if not suspensos_ativos:
                st.info("Nenhum leitor suspenso no momento.")
            else:
                for nome_susp, info_susp in sorted(suspensos_ativos.items()):
                    col_nome, col_remover = st.columns([4, 1])
                    texto_susp = f"- **{nome_susp}** — suspenso(a) até {info_susp['data_fim'].strftime('%d/%m/%Y')}"
                    if info_susp["motivo"]:
                        texto_susp += f" — *{info_susp['motivo']}*"
                    col_nome.markdown(texto_susp)
                    if col_remover.button("Remover", key=f"remover_susp_{nome_susp}"):
                        sh_conn = get_connection()
                        if remover_suspensao(sh_conn, nome_susp):
                            obter_suspensoes.clear()
                            st.success(f"Suspensão de {nome_susp} removida.")
                            time.sleep(2.5)
                            st.rerun()

elif st.session_state.pagina == "quem_nao_serviu":
    st.subheader("Leitores que Ainda Não Serviram")
    if st.session_state.user_profile != "3":
        st.error("Apenas o ADM pode acessar esta tela.")
    else:
        col_sem_ini, col_sem_fim = st.columns(2)
        with col_sem_ini:
            data_sem_inicio = st.date_input("De:", value=date.today().replace(day=1), format="DD/MM/YYYY", key="sem_escala_data_inicio")
        with col_sem_fim:
            data_sem_fim = st.date_input("Até:", value=date.today() + timedelta(days=30), format="DD/MM/YYYY", key="sem_escala_data_fim")

        if data_sem_fim < data_sem_inicio:
            st.error("A data final não pode ser anterior à data inicial.")
        else:
            leitores_que_serviram = set()
            for r in escala_data:
                data_evento_r = extrair_data_evento(str(r.get('DIA', '')))
                if data_evento_r is not None and data_sem_inicio <= data_evento_r <= data_sem_fim:
                    for campo in ('COMENTARISTA', 'LEITURA1', 'LEITURA2'):
                        nome_campo = str(r.get(campo, '')).strip()
                        if nome_campo:
                            leitores_que_serviram.add(nome_campo.strip().upper())

            leitores_sem_escala = [nome for nome in lista_todos_leitores if nome.strip().upper() not in leitores_que_serviram]

            if not leitores_sem_escala:
                st.success("Todos os leitores cadastrados já serviram nesse período!")
            else:
                st.write(f"**{len(leitores_sem_escala)} leitor(es) ainda não serviram entre {data_sem_inicio.strftime('%d/%m/%Y')} e {data_sem_fim.strftime('%d/%m/%Y')}:**")
                for nome_sem_escala in leitores_sem_escala:
                    st.markdown(f"- {nome_sem_escala}")

elif st.session_state.pagina == "escala_geral":
    st.subheader("Escala Geral do Mês")

    col_data_ini, col_data_fim = st.columns(2)
    with col_data_ini:
        data_filtro_inicio = st.date_input("De:", value=date.today(), format="DD/MM/YYYY", key="escala_geral_data_inicio")
    with col_data_fim:
        data_filtro_fim = st.date_input("Até:", value=date.today() + timedelta(days=30), format="DD/MM/YYYY", key="escala_geral_data_fim")

    if data_filtro_fim < data_filtro_inicio:
        st.error("A data final não pode ser anterior à data inicial.")
    else:
        algum_evento = False
        for idx, row in enumerate(escala_data):
            data_evento = extrair_data_evento(str(row.get('DIA', '')))
            if data_evento is not None and not (data_filtro_inicio <= data_evento <= data_filtro_fim):
                continue
            algum_evento = True
            renderizar_evento(idx, row, modo_aguardando=False)

        if not algum_evento:
            st.info("Nenhum evento encontrado nesse intervalo de datas.")

elif st.session_state.pagina == "minha_escala":
    st.subheader("Minha Escala")
    st.info(f"Exibindo eventos agendados para: {st.session_state.user_name}")
    
    encontrou = False
    for idx, row in enumerate(escala_data):
        if not evento_e_hoje_ou_futuro(row):
            continue
        c = str(row.get('COMENTARISTA', '')).strip().upper()
        l1 = str(row.get('LEITURA1', '')).strip().upper()
        l2 = str(row.get('LEITURA2', '')).strip().upper()
        u = st.session_state.user_name.upper()
        
        if c == u or l1 == u or l2 == u:
            encontrou = True
            renderizar_evento(idx, row, modo_aguardando=False)
            
    if not encontrou:
        st.write("Você não possui escalas ativas no momento.")

elif st.session_state.pagina == "exibir_escala":
    st.subheader("Exibir Escala (PDF)")

    col_pdf_ini, col_pdf_fim = st.columns(2)
    with col_pdf_ini:
        pdf_data_inicio = st.date_input("De:", value=date.today(), format="DD/MM/YYYY", key="pdf_escala_data_inicio")
    with col_pdf_fim:
        pdf_data_fim = st.date_input("Até:", value=date.today() + timedelta(days=30), format="DD/MM/YYYY", key="pdf_escala_data_fim")

    if pdf_data_fim < pdf_data_inicio:
        st.error("A data final não pode ser anterior à data inicial.")
    else:
        st.write("Gerando PDF estruturado:")

        class PDF(FPDF):
            def header(self):
                self.set_font('Arial', 'B', 14)
                self.cell(190, 10, 'Escala Geral do Mes', 0, 1, 'C')
                self.ln(5)
            def footer(self):
                self.set_y(-15)
                self.set_font('Arial', 'I', 8)
                self.cell(190, 10, f'Pagina {self.page_no()}', 0, 0, 'C')

        pdf = PDF()
        pdf.add_page()

        algum_evento_pdf = False
        for row in escala_data:
            dia = str(row.get('DIA', ''))
            horario = str(row.get('HORARIO', ''))
            solenidade = str(row.get('SOLENIDADE', 'NÃO')).strip().upper()
            comentarista = str(row.get('COMENTARISTA', '')).strip() or 'Vago'
            l1 = str(row.get('LEITURA1', '')).strip() or 'Vago'
            l2 = str(row.get('LEITURA2', '')).strip() or 'Vago'

            data_evento_pdf = extrair_data_evento(dia)
            if data_evento_pdf is not None and not (pdf_data_inicio <= data_evento_pdf <= pdf_data_fim):
                continue
            algum_evento_pdf = True

            mostrar_comentarista = deve_exibir_comentarista(row)
            mostrar_l2 = deve_exibir_leitura2(row)

            texto_cabecalho = f"Data: {dia} - Horario: {horario}"
            if solenidade == 'SIM':
                texto_cabecalho += " (Solenidade)"
            if eh_dia_4(dia):
                texto_cabecalho += " (Missa da Saude)"

            pdf.set_font("Arial", 'B', 10)
            pdf.cell(190, 7, texto_cabecalho.encode('latin-1', 'replace').decode('latin-1'), 0, 1, 'L')

            pdf.set_font("Arial", '', 10)
            if mostrar_comentarista:
                pdf.cell(190, 6, f"  - COMENTARISTA: {comentarista}".encode('latin-1', 'replace').decode('latin-1'), 0, 1, 'L')
            pdf.cell(190, 6, f"  - 1a LEITURA: {l1}".encode('latin-1', 'replace').decode('latin-1'), 0, 1, 'L')
            if mostrar_l2:
                pdf.cell(190, 6, f"  - 2a LEITURA: {l2}".encode('latin-1', 'replace').decode('latin-1'), 0, 1, 'L')

            pdf.ln(4)

        if not algum_evento_pdf:
            st.info("Nenhum evento encontrado nesse intervalo de datas.")
        else:
            pdf_bytes = bytes(pdf.output())
            b64_pdf = base64.b64encode(pdf_bytes).decode('utf-8')
            st.markdown(f"""
                <a href="data:application/pdf;base64,{b64_pdf}" target="_blank" rel="noopener noreferrer" download="escala_leitores.pdf"
                   style="display:block; text-align:center; background:#0D1B2A; color:#E0E2E5; border:3.5px solid #8C6D4F;
                          border-radius:24px; padding:12px 6px; font-size:18px; font-weight:700; text-decoration:none;
                          margin-top:10px; font-family:sans-serif;">
                    📄 Abrir / Baixar / Compartilhar Escala em PDF
                </a>
            """, unsafe_allow_html=True)
            st.caption("Toque no link para abrir o PDF. A partir da tela de visualização do seu celular, use as opções de baixar, imprimir ou compartilhar.")

elif st.session_state.pagina == "aguardando":
    st.subheader("Aguardando Leitores (Vagas Pendentes)")
    encontrou_vaga = False
    for idx, row in enumerate(escala_data):
        if not evento_e_hoje_ou_futuro(row):
            continue
        mostrar_comentarista = deve_exibir_comentarista(row)
        mostrar_l2 = deve_exibir_leitura2(row)
        c = str(row.get('COMENTARISTA', '')).strip()
        l1 = str(row.get('LEITURA1', '')).strip()
        l2 = str(row.get('LEITURA2', '')).strip()
        
        if (not l1) or (mostrar_comentarista and not c) or (mostrar_l2 and not l2):
            encontrou_vaga = True
            renderizar_evento(idx, row, modo_aguardando=True)
            
    if not encontrou_vaga:
        st.success("Parabéns! Não há vagas pendentes no momento.")

elif st.session_state.pagina == "ver_intencoes":
    st.subheader("Intenções da Santa Missa")
    st.markdown("Selecione abaixo a **Data e o Horário da Missa**:")

    def normalizar_chave(txt):
        txt = str(txt).strip()
        txt = unicodedata.normalize('NFKD', txt).encode('ASCII', 'ignore').decode('ASCII')
        return re.sub(r'\s+', ' ', txt).upper()

    try:
        sh_conn = get_connection()
        ws_resp = sh_conn.worksheet("Respostas")
        respostas_data = ws_resp.get_all_records()

        # Mapa: versão normalizada do cabeçalho -> nome real da coluna na planilha
        # (permite achar a coluna certa mesmo com espaço extra, acento diferente ou maiúsculas/minúsculas)
        colunas_reais = list(respostas_data[0].keys()) if respostas_data else []
        mapa_colunas = {normalizar_chave(c): c for c in colunas_reais}

        def obter_coluna_real(nome_alvo):
            return mapa_colunas.get(normalizar_chave(nome_alvo))

        opcoes_missas = []
        chaves_vistas = set()
        for r in respostas_data:
            d_val = str(r.get('Data', r.get('DATA', ''))).strip()
            h_val = str(r.get('Horário da Missa', r.get('Horario da Missa', r.get('HORARIO', '')))).strip()
            if d_val and h_val:
                chave = (normalizar_chave(d_val), normalizar_chave(h_val))
                if chave not in chaves_vistas:
                    chaves_vistas.add(chave)
                    opcoes_missas.append(f"{d_val} - {h_val}")

        if not opcoes_missas:
            st.info("Nenhuma intenção encontrada na planilha.")
        else:
            missa_selecionada = st.selectbox("Selecione a Missa (Ticket):", sorted(opcoes_missas))

            if st.button("Gerar Relatório Consolidado"):
                partes = missa_selecionada.split(" - ", 1)
                dia_raw = partes[0].strip()
                horario_raw = partes[1].strip() if len(partes) > 1 else ""
                f_data = normalizar_chave(partes[0])
                f_horario = normalizar_chave(partes[1]) if len(partes) > 1 else ""

                registros_encontrados = []
                for r in respostas_data:
                    data_resp = normalizar_chave(r.get('Data', r.get('DATA', '')))
                    horario_resp = normalizar_chave(r.get('Horário da Missa', r.get('Horario da Missa', r.get('HORARIO', ''))))
                    if data_resp == f_data and horario_resp == f_horario:
                        registros_encontrados.append(r)

                if not registros_encontrados:
                    st.info("Nenhum registro encontrado para a missa selecionada.")
                else:
                    st.success(f"{len(registros_encontrados)} envio(s) encontrado(s) e mesclado(s) com sucesso!")
                    time.sleep(2.5)

                    class PDFIntencoes(FPDF):
                        def header(self):
                            if self.page_no() == 1:
                                self.set_font('Arial', 'B', 14)
                                self.cell(190, 10, 'Intenções da Santa Missa'.encode('latin-1', 'replace').decode('latin-1'), 0, 1, 'C')
                                self.set_font('Arial', '', 11)
                                titulo_missa = missa_selecionada.encode('latin-1', 'replace').decode('latin-1')
                                self.cell(190, 8, titulo_missa, 0, 1, 'C')
                                self.ln(3)

                        def footer(self):
                            self.set_y(-15)
                            self.set_font('Arial', 'I', 8)
                            self.cell(190, 10, f'Pagina {self.page_no()}', 0, 0, 'C')

                    pdf = PDFIntencoes()
                    pdf.add_page()

                    # Categorias de intenção: nome de exibição -> nome "alvo" da coluna (resolvido de forma tolerante)
                    categorias_intencao_alvo = [
                        ("Pelas Almas", "Intenções pelas almas"),
                        ("Sétimo Dia", "Missa de sétimo dia"),
                        ("Aniversário Natalício", "Aniversário Natalício"),
                        ("Falecido(a) Hoje", "Falecido(a) Hoje"),
                        ("Bodas", "Bodas"),
                        ("Intenções pela Saúde", "Intenções pela Saúde"),
                    ]
                    categorias_intencao = [
                        (titulo, obter_coluna_real(alvo) or alvo)
                        for titulo, alvo in categorias_intencao_alvo
                    ]
                    mapa_categorias = dict(categorias_intencao)

                    def coletar_linhas(coluna):
                        linhas = []
                        for r in registros_encontrados:
                            valor = str(r.get(coluna, '')).strip()
                            if valor:
                                for linha_texto in valor.splitlines():
                                    linha_texto = linha_texto.strip()
                                    if linha_texto:
                                        linhas.append(linha_texto)
                        return linhas

                    def preencher_caixa(x, y_topo, y_fundo, largura, titulo, linhas, repetir_titulo=True):
                        """Escreve um título + lista de linhas dentro de uma caixa (x, y_topo)-(x+largura, y_fundo).
                        Para quando o espaço acaba e devolve as linhas que não couberam."""
                        pdf.set_xy(x, y_topo)
                        if repetir_titulo:
                            pdf.set_font("Arial", 'B', 11)
                            pdf.set_x(x)
                            pdf.cell(largura, 6, titulo.encode('latin-1', 'replace').decode('latin-1'), 0, 2, 'L')

                        if not linhas:
                            pdf.set_font("Arial", 'I', 9)
                            pdf.set_x(x)
                            pdf.multi_cell(largura, 5, "(nenhuma intenção informada)".encode('latin-1', 'replace').decode('latin-1'))
                            return []

                        pdf.set_font("Arial", '', 9)
                        for i, linha_texto in enumerate(linhas):
                            if pdf.get_y() + 5 > y_fundo:
                                return linhas[i:]
                            pdf.set_x(x)
                            pdf.multi_cell(largura, 5, linha_texto.encode('latin-1', 'replace').decode('latin-1'))
                        return []

                    # --- Layout da página 1: duas colunas ---
                    Y0 = pdf.get_y()
                    ALTURA_TOTAL = 250.0
                    X_ESQ, LARG_ESQ = 10, 90
                    X_DIR, LARG_DIR = 105, 95

                    y_esq_almas_fim = Y0 + ALTURA_TOTAL * 0.80
                    y_esq_aniversario_fim = Y0 + ALTURA_TOTAL * 1.00

                    y_dir_almas_cont_fim = Y0 + ALTURA_TOTAL * 0.70
                    y_dir_falecido_fim = y_dir_almas_cont_fim + ALTURA_TOTAL * 0.10
                    y_dir_setimo_fim = y_dir_falecido_fim + ALTURA_TOTAL * 0.10
                    y_dir_bodas_fim = y_dir_setimo_fim + ALTURA_TOTAL * 0.10

                    linhas_almas = coletar_linhas(mapa_categorias["Pelas Almas"])
                    sobra_almas = preencher_caixa(X_ESQ, Y0, y_esq_almas_fim, LARG_ESQ, "Intenções pelas Almas", linhas_almas)

                    linhas_aniversario = coletar_linhas(mapa_categorias["Aniversário Natalício"])
                    preencher_caixa(X_ESQ, y_esq_almas_fim, y_esq_aniversario_fim, LARG_ESQ, "Aniversário Natalício", linhas_aniversario)

                    if sobra_almas:
                        preencher_caixa(X_DIR, Y0, y_dir_almas_cont_fim, LARG_DIR, "Intenções pelas Almas (continuação)", sobra_almas)

                    linhas_falecido = coletar_linhas(mapa_categorias["Falecido(a) Hoje"])
                    preencher_caixa(X_DIR, y_dir_almas_cont_fim, y_dir_falecido_fim, LARG_DIR, "Falecido(a) Hoje", linhas_falecido)

                    linhas_setimo = coletar_linhas(mapa_categorias["Sétimo Dia"])
                    preencher_caixa(X_DIR, y_dir_falecido_fim, y_dir_setimo_fim, LARG_DIR, "Missa de Sétimo Dia", linhas_setimo)

                    linhas_bodas = coletar_linhas(mapa_categorias["Bodas"])
                    preencher_caixa(X_DIR, y_dir_setimo_fim, y_dir_bodas_fim, LARG_DIR, "Bodas", linhas_bodas)

                    # --- Página 2: Intenções pela Saúde ---
                    pdf.add_page()
                    linhas_saude = coletar_linhas(mapa_categorias["Intenções pela Saúde"])
                    pdf.set_font("Arial", 'B', 14)
                    pdf.set_xy(10, pdf.get_y())
                    pdf.cell(190, 8, "Intenções pela Saúde".encode('latin-1', 'replace').decode('latin-1'), 0, 2, 'C')
                    pdf.ln(2)
                    preencher_caixa(10, pdf.get_y(), 280, 190, "", linhas_saude, repetir_titulo=False)

                    pdf_bytes = bytes(pdf.output())
                    b64_pdf = base64.b64encode(pdf_bytes).decode('utf-8')
                    st.markdown(f"""
                        <a href="data:application/pdf;base64,{b64_pdf}" target="_blank" rel="noopener noreferrer" download="Intenções_da_Santa_Missa.pdf"
                           style="display:block; text-align:center; background:#0D1B2A; color:#FFFFFF; border:3.5px solid #8C6D4F;
                                  border-radius:24px; padding:12px 6px; font-size:18px; font-weight:700; text-decoration:none;
                                  margin-top:10px; font-family:sans-serif;">
                            📄 Abrir / Baixar / Compartilhar Intenções em PDF
                        </a>
                    """, unsafe_allow_html=True)
                    st.caption("Toque no link para abrir o PDF. A partir da tela de visualização do seu celular, use as opções de baixar, imprimir ou compartilhar.")

                    # --- Exportação para apresentação (.pptx), um slide por categoria ---
                    from pptx import Presentation
                    from pptx.util import Inches, Pt
                    from pptx.dml.color import RGBColor
                    import io

                    prs = Presentation()
                    prs.slide_width = Inches(13.333)
                    prs.slide_height = Inches(7.5)
                    layout_em_branco = prs.slide_layouts[6]

                    COR_FUNDO = RGBColor(0xFE, 0xFA, 0xE0)
                    COR_TITULO = RGBColor(0x5C, 0x3A, 0x21)
                    COR_TEXTO = RGBColor(0x3D, 0x26, 0x12)

                    for titulo_categoria, coluna in categorias_intencao:
                        slide = prs.slides.add_slide(layout_em_branco)
                        slide.background.fill.solid()
                        slide.background.fill.fore_color.rgb = COR_FUNDO

                        caixa_titulo = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.2))
                        tf_titulo = caixa_titulo.text_frame
                        tf_titulo.text = titulo_categoria
                        tf_titulo.paragraphs[0].font.size = Pt(40)
                        tf_titulo.paragraphs[0].font.bold = True
                        tf_titulo.paragraphs[0].font.color.rgb = COR_TITULO

                        linhas_categoria = []
                        for r in registros_encontrados:
                            valor = str(r.get(coluna, '')).strip()
                            if valor:
                                for linha_texto in valor.splitlines():
                                    linha_texto = linha_texto.strip()
                                    if linha_texto:
                                        linhas_categoria.append(linha_texto)

                        caixa_corpo = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.2))
                        tf_corpo = caixa_corpo.text_frame
                        tf_corpo.word_wrap = True

                        if not linhas_categoria:
                            tf_corpo.text = "(nenhuma intenção informada)"
                            tf_corpo.paragraphs[0].font.italic = True
                            tf_corpo.paragraphs[0].font.size = Pt(24)
                            tf_corpo.paragraphs[0].font.color.rgb = COR_TEXTO
                        else:
                            tf_corpo.text = linhas_categoria[0]
                            tf_corpo.paragraphs[0].font.size = Pt(24)
                            tf_corpo.paragraphs[0].font.color.rgb = COR_TEXTO
                            for linha_texto in linhas_categoria[1:]:
                                p_linha = tf_corpo.add_paragraph()
                                p_linha.text = linha_texto
                                p_linha.font.size = Pt(24)
                                p_linha.font.color.rgb = COR_TEXTO

                    pptx_buffer = io.BytesIO()
                    prs.save(pptx_buffer)
                    b64_pptx = base64.b64encode(pptx_buffer.getvalue()).decode('utf-8')

                    nome_arquivo_pptx = f"intenções_missa{dia_raw.replace('/', '')}_{horario_raw.replace(':', 'h')}.pptx"

                    st.markdown(f"""
                        <a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_pptx}" download="{nome_arquivo_pptx}"
                           style="display:block; text-align:center; background:#0D1B2A; color:#FFFFFF; border:3.5px solid #8C6D4F;
                                  border-radius:24px; padding:12px 6px; font-size:18px; font-weight:700; text-decoration:none;
                                  margin-top:10px; font-family:sans-serif;">
                            🖥️ Baixar Apresentação (PPTX)
                        </a>
                    """, unsafe_allow_html=True)
                    st.caption("Um slide por categoria de intenção, já mesclando todos os envios encontrados.")
    except Exception as e:
        st.error(f"Erro ao acessar a aba de respostas: {e}")